from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Monetstudy"
    subtitle.text = "The intelligence to learn your way.\n\nPresenter: Kyeyune Neithen, Founder"

    def add_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.shapes.placeholders[1].text_frame
        tf.word_wrap = True
        
        # Add first bullet
        if isinstance(bullet_points[0], tuple):
            p = tf.paragraphs[0]
            p.text = bullet_points[0][0]
            p.level = bullet_points[0][1]
        else:
            tf.text = bullet_points[0]
            
        # Add remaining bullets
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            if isinstance(point, tuple):
                p.text = point[0]
                p.level = point[1]
            else:
                p.text = point
                p.level = 0

    # Slide 2: The Problem
    add_slide("The Problem", [
        "Information Overload: Students are overwhelmed by disorganized notes and scattered resources.",
        "Dense, Intimidating Textbooks: Massive books discourage students.",
        "Standardized Inefficiency: Traditional education forces everyone to learn the same way."
    ])

    # Slide 3: The Philosophy
    add_slide("The Philosophy", [
        "The Core Belief: The best way to learn is our own way.",
        "The Shift: Education shouldn't force the student to adapt to the material. It should focus on adapting the material to the student."
    ])

    # Slide 4: The Solution
    add_slide("The Solution: Monetstudy", [
        "A next-generation AI-powered personalized learning environment.",
        "Students upload hundreds of pages of unorganized material or dense textbooks.",
        "They apply Custom Instructions (e.g., 'Teach me using real-world analogies').",
        "Monetstudy instantly transforms the chaos into a structured, perfectly tailored course."
    ])

    # Slide 5: How It Works
    add_slide("How It Works", [
        "Step 1: Ingest. Upload massive documents. Our AI pipeline handles entire textbooks seamlessly.",
        "Step 2: Personalize. Set Custom Instructions. Tell the AI how your brain processes info.",
        "Step 3: Learn. Get a dynamic course with beautifully rendered equations, rich AI illustrations, and step-by-step topic breakdowns."
    ])

    # Slide 6: The Unfair Advantage (Technology)
    add_slide("The Unfair Advantage (Technology)", [
        "Massive Document Processing: Custom-built architecture deeply analyzes and retrieves data from 3M+ character documents natively.",
        "Intelligent Generation: Acts as a dynamic tutor creating topical outlines.",
        "Integrated AI Visuals: Embeds custom visual illustrations directly alongside the text to aid memory."
    ])

    # Slide 7: Business Model & Pricing
    add_slide("Business Model & Pricing", [
        "Freemium Strategy to remove barriers while monetizing heavy users.",
        ("Tier 1: Free / Starter - Core course generation.", 1),
        ("Tier 2: Scholar Plan ($5/month) - Premium features including AI-generated illustrations.", 1),
        ("Tier 3: Unlimited Plan ($10/month) - Full unrestricted processing for massive databases.", 1),
        "Seamless Payments: Natively integrated with Pesapal for mobile money and local cards to ensure accessibility across Africa."
    ])

    # Slide 8: The Market Opportunity
    add_slide("The Market Opportunity", [
        "Local to Global: A massive, underserved demographic of students locked out of expensive private tutoring.",
        "The Ed-Tech Gap: Monetstudy democratizes elite, personalized tutoring.",
        "We are providing a world-class study tool that scales infinitely, all for the price of a cup of coffee."
    ])

    # Slide 9: The Team
    add_slide("The Team", [
        "Kyeyune Neithen - Founder & Lead Engineer",
        ("Solo architected the full-stack web application from a room in Uganda.", 1),
        ("Custom-built payment integrations, advanced AI document processing, and cloud architecture.", 1),
        ("Deeply passionate about solving the global education crisis through personalized technology.", 1)
    ])

    # Slide 10: The Ask / Vision
    add_slide("The Ask / Vision", [
        "The Vision: To become the default study layer for the modern student.",
        "Turning any disorganized notes or dense textbook into a highly personalized, interactive learning journey.",
        "Call to Action: Looking for initial seed funding or strategic partnerships to scale server infrastructure and accelerate user acquisition."
    ])

    prs.save('Monetstudy_Pitch_Deck.pptx')
    print("PowerPoint presentation generated successfully!")

if __name__ == '__main__':
    create_presentation()
